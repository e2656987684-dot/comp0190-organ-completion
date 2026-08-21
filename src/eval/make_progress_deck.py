"""Build the second progress-report deck (phase starting 2026-08-07).

Deliberately light: seven slides, text-led, two small tables, no figures. The
first report already carried the density pictures, and this phase's material is
mostly results and decisions rather than new visuals.

Scope was set by the student:
  * cover what changed since the 2026-08-07 ablation
  * the results slides may be technical -- those are the ones being presented
  * the attention-behaviour slide must be readable at a glance: short bullets,
    markers, no paragraphs. It is the one part that is hard to field questions on
  * two slides of thesis outline, with a working title
  * no reimplementation-audit slide, and no limitations slide built out of work
    not yet done

    python src/eval/make_progress_deck.py
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
OUT = os.path.join(REPO, "reports", "progress_report_2.pptx")

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5B, 0x66, 0x73)
ACCENT = RGBColor(0x15, 0x65, 0xC0)
RULE = RGBColor(0xD5, 0xDA, 0xE0)


def _text(slide, x, y, w, h, runs, size=16, color=INK, bold=False,
          align=PP_ALIGN.LEFT, space=10):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (txt, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space", space))
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
        r.font.name = "Calibri"
    return box


def _slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _text(s, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.7),
          title, size=28, bold=True)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.05), Inches(12.2), Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return s


def _table(slide, x, y, w, h, rows, col_w=None, size=14, highlight=None):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(0.36)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = para.runs[0] if para.runs else para.add_run()
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            hot = highlight is not None and (ri, ci) == highlight
            r.font.bold = (ri == 0) or (ci == 0) or hot
            r.font.color.rgb = ACCENT if hot else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xEF, 0xF3, 0xF7) if ri == 0
                                        else RGBColor(0xFF, 0xFF, 0xFF))
    return tbl


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    L, TOP, COL = Inches(0.7), Inches(1.45), Inches(11.9)

    # ---------------------------------------------------------------- 1
    s = _slide(prs, "Where things stand")
    _text(s, L, TOP, Inches(6.4), Inches(4.6), [
        ("The task", {"bold": True, "size": 17, "space": 6}),
        ("4,096 points of a defective skull in, 6,144 points of a complete skull out. "
         "SkullFix, 80 train / 20 validation, the same split throughout so runs stay "
         "comparable.", {"space": 16}),
        ("Reported in millimetres", {"bold": True, "size": 17, "space": 6}),
        ("Every cloud carries its own scale factor, so Chamfer and Hausdorff are converted "
         "back to mm rather than left in normalised units.", {}),
    ], size=15.5)
    _text(s, Inches(7.3), TOP, Inches(5.3), Inches(4.6), [
        ("Since the last report", {"bold": True, "size": 17, "space": 10}),
        ("1.  Completed the loss ablation that was left open", {"space": 9}),
        ("2.  Added metrics restricted to the defect region", {"space": 9}),
        ("3.  Measured which parts of the network are actually used", {"space": 9}),
        ("4.  Started the thesis, and started cross-validation", {"space": 9}),
    ], size=15.5)

    # ---------------------------------------------------------------- 2
    s = _slide(prs, "Loss ablation: is DCD still needed?")
    _text(s, L, TOP, COL, Inches(1.5), [
        ("Repulsion is a hinge penalty on each point's nearest neighbours — it pushes any "
         "pair closer than 2 mm apart, and is normalised by that radius so the weight "
         "means the same thing at any scale. It fixed the clumping, which left the "
         "question: with repulsion in place, does the paper's density-aware loss (DCD) "
         "still contribute?", {"size": 15.5}),
    ])
    _table(s, Inches(1.0), Inches(3.05), Inches(11.3), Inches(1.45), [
        ["Chamfer distance / points closer than 2 mm", "no repulsion", "with repulsion"],
        ["with DCD", "6.40 mm  /  5.6 %", "6.41 mm  /  1.4 %"],
        ["no DCD", "6.43 mm  /  13.6 %", "6.36 mm  /  1.3 %"],
    ], col_w=[Inches(4.7), Inches(3.3), Inches(3.3)], highlight=(2, 2))
    _text(s, L, Inches(4.85), COL, Inches(2.2), [
        ("Chamfer + repulsion wins, and it is the simplest of the four.", {"bold": True, "space": 8}),
        ("•  Repeating a run with identical settings moves Chamfer distance by 0.004 mm, "
         "so the 0.07 mm spread across this table is real and not noise.", {"space": 7, "size": 15}),
        ("•  DCD alone does work — clumping 13.6 % → 5.6 % — but repulsion reaches 1.3 % "
         "on its own, and stacking the two gains nothing.", {"space": 7, "size": 15}),
        ("•  Accuracy is flat across all four. Density and accuracy are being controlled "
         "by different things.", {"space": 7, "size": 15}),
    ])

    # ---------------------------------------------------------------- 3
    s = _slide(prs, "Scoring inside the defect")
    _text(s, L, TOP, COL, Inches(1.5), [
        ("A whole-cloud score is dominated by the intact skull, which the model can copy "
         "from its input. To score the part that matters, a ground-truth point counts as "
         "\"in the defect\" when the nearest input point is more than 5 mm away. That "
         "threshold sits in the valley of a clearly bimodal distance distribution, and "
         "selects about 6 % of the ground-truth points.", {"size": 15.5}),
    ])
    _table(s, Inches(1.9), Inches(3.15), Inches(9.5), Inches(1.8), [
        ["configuration", "defect coverage", "defect precision"],
        ["Chamfer only", "3.44 mm", "2.93 mm"],
        ["Chamfer + DCD", "3.25 mm", "2.89 mm"],
        ["Chamfer + DCD + repulsion", "3.41 mm", "2.96 mm"],
        ["Chamfer + repulsion", "3.24 mm", "2.91 mm"],
    ], col_w=[Inches(4.3), Inches(2.6), Inches(2.6)], highlight=(4, 1))
    _text(s, L, Inches(5.3), COL, Inches(1.8), [
        ("•  Improvements are larger inside the hole (−17 %) than over the whole skull "
         "(−12 %), so the changes are helping where the problem is hard.", {"space": 7, "size": 15}),
        ("•  Coverage separates the configurations; precision sits at 2.89–2.96 mm for all "
         "of them and tells you nothing. Coverage is the metric to report.",
         {"space": 7, "size": 15}),
    ])

    # ---------------------------------------------------------------- 4
    s = _slide(prs, "A look inside the model")
    _text(s, L, TOP, COL, Inches(0.6),
          "Does the attention actually do anything?", size=19, bold=True)
    _text(s, L, Inches(2.25), Inches(6.0), Inches(4.4), [
        ("16", {"bold": True, "size": 40, "color": ACCENT, "space": 0}),
        ("attention blocks in the network", {"size": 15, "color": MUTED, "space": 22}),
        ("✗   12 of them  →  every point gets", {"size": 17, "space": 2}),
        ("      the same weight", {"size": 17, "space": 6}),
        ("      = an averaging step, not attention", {"size": 14, "color": MUTED, "space": 16}),
        ("✓   4 of them  →  really do focus", {"size": 17, "space": 2}),
        ("      on particular points", {"size": 17, "space": 0}),
    ])
    _text(s, Inches(7.1), Inches(2.25), Inches(5.5), Inches(4.4), [
        ("So what?", {"bold": True, "size": 19, "space": 14}),
        ("→   I tried adding more attention.", {"size": 16, "space": 3}),
        ("      Results got worse.", {"size": 16, "bold": True, "space": 16}),
        ("→   Now I know why, so I am not", {"size": 16, "space": 3}),
        ("      spending more time there.", {"size": 16, "space": 16}),
        ("→   Effort goes to the loss functions", {"size": 16, "space": 3}),
        ("      and the evaluation instead.", {"size": 16, "space": 0}),
    ])

    # ---------------------------------------------------------------- 5
    s = _slide(prs, "Thesis outline  (1 / 2)")
    _text(s, L, Inches(1.35), COL, Inches(0.9), [
        ("Point-Cloud Completion of Cranial Defects:", {"bold": True, "size": 21, "space": 2}),
        ("Density-Aware Losses and Defect-Region Evaluation", {"bold": True, "size": 21,
                                                               "color": ACCENT, "space": 0}),
    ])
    _text(s, L, Inches(2.75), COL, Inches(4.2), [
        ("1.  Introduction", {"bold": True, "size": 17, "space": 3}),
        ("Cranial implant design and why it is done by hand today. Why a point cloud "
         "rather than a voxel grid, and what that choice costs.", {"size": 15, "space": 16}),
        ("2.  Background", {"bold": True, "size": 17, "space": 3}),
        ("Point-cloud completion from PCN onwards, transformer-based methods, and the "
         "multimodal model this work builds on. The datasets in this area and their sizes.",
         {"size": 15, "space": 16}),
        ("3.  Method", {"bold": True, "size": 17, "space": 3}),
        ("The architecture as reimplemented here, the data pipeline from meshes to aligned "
         "point clouds, and how the training setup differs from the published one.",
         {"size": 15, "space": 0}),
    ])

    # ---------------------------------------------------------------- 6
    s = _slide(prs, "Thesis outline  (2 / 2)")
    _text(s, L, Inches(1.35), COL, Inches(5.4), [
        ("4.  Evaluation protocol", {"bold": True, "size": 17, "space": 3}),
        ("The metrics used and why. The defect-region metrics and how the region is "
         "defined. What the point-cloud representation itself limits, and why the numbers "
         "here cannot be placed next to published voxel-domain ones.", {"size": 15, "space": 13}),
        ("5.  Loss functions", {"bold": True, "size": 17, "space": 3}),
        ("The point-density problem, the repulsion term, and the ablation that settles "
         "which loss terms are needed and which can be dropped.", {"size": 15, "space": 13}),
        ("6.  Model behaviour", {"bold": True, "size": 17, "space": 3}),
        ("A short chapter on which parts of the network are actually being used, and what "
         "that implies for where effort is worth spending.", {"size": 15, "space": 13}),
        ("7.  Discussion", {"bold": True, "size": 17, "space": 3}),
        ("What the results say about this class of model on data of this size, and how the "
         "findings would transfer.", {"size": 15, "space": 13}),
        ("8.  Conclusion", {"bold": True, "size": 17, "space": 0}),
    ])

    # ---------------------------------------------------------------- 7
    s = _slide(prs, "Next steps")
    _text(s, L, TOP, COL, Inches(4.6), [
        ("Cross-validation is running.", {"bold": True, "size": 17, "space": 4}),
        ("Longer than a single run would take: the ablation left several configurations to "
         "cover, and I have moved to a larger dataset.", {"space": 16}),
        ("Writing has started.", {"bold": True, "size": 17, "space": 4}),
        ("On the parts that are already settled and will not change.", {"space": 16}),
        ("Then: a surface-reconstruction comparison.", {"bold": True, "size": 17, "space": 4}),
        ("So the point-cloud numbers can be related to the ones clinical work is usually "
         "reported in.", {"space": 0}),
    ], size=15.5)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"-> {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
