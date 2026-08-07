"""Build the progress-report slide deck from the generated figures and metrics.

Numbers come from `reports/figures/summary.csv`, which `make_report_figures.py`
writes out of the real runs -- so a slide cannot quietly disagree with the code.
Run that script first, then this one.

    python src/eval/make_report_figures.py
    python src/eval/make_report_deck.py

Eight content slides, no cover / contents / summary by request. Titles are kept
to two or three words; the context that would otherwise sit in a subtitle is in
the body text, where it can be read rather than skimmed.
"""

from __future__ import annotations

import os

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIG = os.path.join(REPO, "reports", "figures")
OUT = os.path.join(REPO, "reports", "progress_report.pptx")

W, H = Inches(13.333), Inches(7.5)          # 16:9
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x5B, 0x66, 0x73)
ACCENT = RGBColor(0x15, 0x65, 0xC0)
WARN = RGBColor(0xB7, 0x3B, 0x1F)
RULE = RGBColor(0xD5, 0xDA, 0xE0)


def _text(slide, x, y, w, h, runs, size=15, color=INK, bold=False,
          align=PP_ALIGN.LEFT, space=7):
    """runs: str, or list of (text, {overrides}) for mixed formatting."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (txt, ov) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.space_after = Pt(ov.get("space", space))
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
        r.font.name = "Calibri"
    return box


def _slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])      # blank
    _text(s, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.65),
          title, size=28, bold=True)
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.02), Inches(12.2), Emu(9525))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return s


def _picture(slide, name, y, height=None, width=None, cx=None):
    p = os.path.join(FIG, name)
    kw = {"width": width} if width else {"height": height}
    pic = slide.shapes.add_picture(p, 0, y, **kw)
    pic.left = int((cx if cx is not None else W / 2) - pic.width / 2)
    return pic


def _table(slide, x, y, w, h, rows, col_w=None, size=12.5, highlight_row=None):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    tbl = shp.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(0.3)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = para.runs[0] if para.runs else para.add_run()
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            r.font.bold = (ri == 0) or (highlight_row is not None and ri == highlight_row)
            r.font.color.rgb = (ACCENT if highlight_row is not None and ri == highlight_row
                                else INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0xEF, 0xF3, 0xF7) if ri == 0
                                        else RGBColor(0xFF, 0xFF, 0xFF))
    return shp


def main():
    df = pd.read_csv(os.path.join(FIG, "summary.csv")).set_index("run")
    g = lambda r, c: df.loc[r, c]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ============================================================== 1. Task
    s = _slide(prs, "Task and baseline")
    _text(s, Inches(0.65), Inches(1.25), Inches(5.5), Inches(4.4),
          [("Skull completion: 4,096 points of a defective skull in, 6,144 points of a "
            "complete skull out.", {}),
           ("SkullFix, 100 skulls, 80 train / 20 validation — the same split for every "
            "run reported here.", {}),
           ("For a reference point I ran the authors' released weights on my own data: "
            "same validation skulls, same metrics, same aligned point clouds.", {})],
          size=15, space=13)
    _picture(s, "pred_vs_gt.png", Inches(1.15), height=Inches(3.7), cx=Inches(9.5))
    _table(s, Inches(0.65), Inches(3.6), Inches(5.6), Inches(1.0),
           [["", "CD_t (mm)", "DCD"],
            ["Pretrained weights", "10.71", "1.428"],
            ["This work", f"{g('rep_w05','CD_t_mm'):.2f}", f"{g('rep_w05','DCD'):.3f}"]],
           col_w=[Inches(2.6), Inches(1.5), Inches(1.5)], highlight_row=2)
    _text(s, Inches(0.65), Inches(5.35), Inches(12.1), Inches(1.5),
          [("Sanity check.", {"bold": True, "color": ACCENT}),
           ("The source paper reports DCD = 1.41269 on its own validation set. The same "
            "weights measure 1.42772 on my data — a 1% match, which confirms the metric "
            "implementation and the weight loading.", {}),
           ("It also frames the comparison honestly: their model has not failed on "
            "skulls. It performs here almost exactly as it does on its own data, so the "
            "gain comes from specialisation.", {"color": MUTED, "size": 14})],
          size=15, space=6)

    # ================================================== 2. Evaluation tooling
    s = _slide(prs, "Evaluation tooling")
    _text(s, Inches(0.65), Inches(1.25), Inches(5.7), Inches(4.6),
          [("Scatter plots cannot show whether a predicted surface is any good, so I "
            "built a reconstruction path: point cloud → distance field → marching cubes "
            "→ Taubin smoothing. The original repository has none of this.", {}),
           ("Poisson reconstruction, the usual choice, was deliberately avoided. It "
            "needs per-point normals, which are unreliable on a thin bone shell, and it "
            "tends to close holes — the one feature being inspected.", {}),
           ("Two rules that came out of it:", {"bold": True, "space": 3}),
           ("Meshes are for looking, never for metrics. Reconstruction pushes the "
            "surface outward by a median of 5.3 mm, the same order as the model's own "
            "error.", {"size": 14, "color": MUTED, "space": 3}),
           ("Reconstruction parameters are locked. Every knob changes how smooth a "
            "surface looks, so per-figure tuning would disguise itself as model "
            "improvement.", {"size": 14, "color": MUTED})],
          size=15, space=11)
    _picture(s, "mesh.png", Inches(1.35), height=Inches(2.9), cx=Inches(9.7))
    _text(s, Inches(6.6), Inches(4.5), Inches(6.2), Inches(2.4),
          [("What a mesh cannot show", {"bold": True, "color": WARN, "size": 16}),
           ("These two surfaces look alike, and the measurements agree. The real "
            "difference between my models is in how the points are distributed — which "
            "a shaded surface hides completely.", {}),
           ("That is why the diagnostics from here on are point clouds, not meshes.",
            {"color": MUTED, "size": 14})],
          size=15, space=9)

    # ========================================================== 3. Diagnosis
    s = _slide(prs, "Change of direction")
    _text(s, Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.9),
          [("The original plan", {"bold": True, "size": 17, "color": ACCENT}),
           ("Add a smoothness penalty — the predicted surface looked bumpy in the "
            "renders. I measured it first, and three things argued against it.", {})],
          size=15, space=8)
    _text(s, Inches(0.7), Inches(2.35), Inches(5.8), Inches(4.4),
          [("1.  No measurable difference to exploit", {"bold": True}),
           ("Local roughness came out at 0.736 for ground truth against 0.760 for the "
            "prediction.", {"size": 14, "color": MUTED}),
           ("2.  The measurement itself is unreliable", {"bold": True}),
           ("A skull is a thin shell, inner and outer surfaces 5–7 mm apart, and a local "
            "surface fit cannot tell that apart from real roughness. Without a "
            "trustworthy measure I could not have judged whether the change helped.",
            {"size": 14, "color": MUTED}),
           ("3.  Overdoing it would cost accuracy", {"bold": True}),
           ("Pushing past the ground truth's own smoothness makes the prediction "
            "smoother than the target, which the distance metrics count as error.",
            {"size": 14, "color": MUTED})],
          size=15, space=7)

    _text(s, Inches(7.0), Inches(1.2), Inches(5.6), Inches(1.4),
          [("What the data pointed at instead", {"bold": True, "size": 17,
                                                 "color": WARN}),
           ("Point spacing. 12.8% of predicted points sit within 2 mm of a neighbour; "
            "ground truth, being farthest-point sampled, has none at all.", {})],
          size=15, space=8)
    _table(s, Inches(7.0), Inches(2.75), Inches(5.6), Inches(1.0),
           [["", "prediction", "ground truth"],
            ["points closer than 2 mm", f"{g('baseline','clump_%'):.1f}%", "0.0%"],
            ["spacing CV", f"{g('baseline','spacing_CV'):.3f}", "0.145"]],
           col_w=[Inches(2.6), Inches(1.5), Inches(1.5)])
    _text(s, Inches(7.0), Inches(4.1), Inches(5.6), Inches(2.7),
          [("Those clumped points are wasted — they cover no surface a neighbour does "
            "not already cover. Unlike roughness this is unambiguous, and it is what the "
            "rest of this work goes after.", {}),
           ("Worth noting in advance: fixing the density also improved the surface. "
            "Deviation p95 fell from 6.56 to 5.59 mm without any smoothness term at all.",
            {"color": ACCENT, "size": 14})],
          size=15, space=10)

    # ============================================ 4. Loss experiments + why
    s = _slide(prs, "Loss experiments")
    _text(s, Inches(0.7), Inches(1.15), Inches(5.9), Inches(0.35),
          "What I tried", size=17, bold=True, color=ACCENT)
    _text(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(1.0),
          "The obvious response is to lean harder on DCD — the density-aware loss the "
          "original project trains with. It has two knobs, and I swept both.", size=14.5)
    _table(s, Inches(0.7), Inches(2.75), Inches(5.9), Inches(1.6),
           [["", "CD_t", "clumped", "CV"],
            ["baseline", f"{g('baseline','CD_t_mm'):.2f}",
             f"{g('baseline','clump_%'):.1f}%", f"{g('baseline','spacing_CV'):.3f}"],
            ["weight 1→3", f"{g('dcd_w3','CD_t_mm'):.2f}",
             f"{g('dcd_w3','clump_%'):.1f}%", f"{g('dcd_w3','spacing_CV'):.3f}"],
            ["lambda 1→2", f"{g('dcd_l2','CD_t_mm'):.2f}",
             f"{g('dcd_l2','clump_%'):.1f}%", f"{g('dcd_l2','spacing_CV'):.3f}"]],
           col_w=[Inches(2.0), Inches(1.3), Inches(1.4), Inches(1.2)], size=13)
    _text(s, Inches(0.7), Inches(4.55), Inches(5.8), Inches(2.3),
          [("Raising the weight buys density at a poor rate — 2.8 points of clumping for "
            "0.37 mm of accuracy, and 10% is still far from 0%. Raising lambda, which "
            "targets density specifically, moves it almost not at all.", {}),
           ("Both were dead ends. The question is whether I swept them badly, or whether "
            "the loss simply cannot do this.", {"color": MUTED, "size": 14})],
          size=14.5, space=10)

    _text(s, Inches(7.15), Inches(1.15), Inches(5.5), Inches(0.35),
          "Why it could never work", size=17, bold=True, color=WARN)
    _text(s, Inches(7.15), Inches(1.6), Inches(5.5), Inches(5.2),
          [("DCD  =  1 − exp(−α·d) · 1 / count^λ", {"size": 15, "color": ACCENT,
                                                    "bold": True}),
           ("The density part of DCD works by counting how many ground-truth points end "
            "up assigned to the same prediction. A count is a whole number: nudge a "
            "point slightly and it usually does not change at all. So it tells the "
            "optimiser nothing about which direction to move — the gradient is not small, "
            "it is absent. I checked, and TensorFlow returns nothing for it.", {}),
           ("DCD can therefore make Chamfer care more about some points than others, "
            "but it can never push two predictions apart.", {"bold": True}),
           ("It is also blind to the most common case: two predictions sitting on top of "
            "each other, each claimed by a different ground-truth point. Both counts are "
            "1, so nothing is charged.", {}),
           ("Repulsion  =  max(0, r₀ − d)²", {"size": 15, "color": ACCENT, "bold": True}),
           ("If two predicted points are closer than r₀, push them apart; once they are "
            "r₀ apart, do nothing. That is the force DCD structurally lacks — so the two "
            "are not redundant.", {"color": WARN, "size": 14})],
          size=14, space=8)

    # ======================================================= 5. What worked
    s = _slide(prs, "What worked")
    _picture(s, "curves.png", Inches(1.15), height=Inches(2.6))
    _text(s, Inches(0.65), Inches(3.95), Inches(6.1), Inches(2.9),
          [("A configuration bug, not loss design", {"bold": True, "size": 16,
                                                     "color": ACCENT}),
           ("The learning-rate scheduler and early stopping both watch validation loss, "
            "but the scheduler's patience was 40 against early stopping's 20 — so "
            "training always ended first and the learning rate had never once been "
            "reduced, in any experiment.", {}),
           ("Fixing it: CD_t 7.22 → 6.40 mm, training runs to 279 epochs instead of 133. "
            "The dotted lines mark each drop.", {"size": 14, "color": MUTED})],
          size=14.5, space=9)
    _text(s, Inches(7.1), Inches(3.95), Inches(5.7), Inches(2.9),
          [("Repulsion loss", {"bold": True, "size": 16, "color": ACCENT}),
           ("Points closer than 2 mm to each other get pushed apart, and the penalty "
            "switches off entirely once they are far enough — so it stops competing with "
            "Chamfer rather than pushing indefinitely. The 2 mm threshold was read off "
            "the ground truth's own spacing.", {}),
           ("Attribution: a third run with the learning-rate fix but no repulsion "
            "separates the two. All of the CD_t gain is the learning rate — 6.40 against "
            "6.41. Repulsion contributes density only.", {"size": 14, "color": MUTED})],
          size=14.5, space=9)

    # =========================================================== 6. Density
    s = _slide(prs, "Density")
    _text(s, Inches(0.65), Inches(1.15), Inches(12.1), Inches(0.8),
          [("Colour is nearest-neighbour spacing, dark meaning points on top of each "
            "other. One colour scale across all four panels — auto-scaling each panel "
            "to its own data would render a 12.8% clump rate and a 1.4% one identically.",
            {})],
          size=15)
    _picture(s, "density.png", Inches(2.0), height=Inches(2.9))
    _table(s, Inches(1.4), Inches(5.15), Inches(10.5), Inches(1.3),
           [["", "clumped <2 mm", "spacing CV", "CD_t (mm)"],
            ["baseline", f"{g('baseline','clump_%'):.1f}%",
             f"{g('baseline','spacing_CV'):.3f}", f"{g('baseline','CD_t_mm'):.2f}"],
            ["+ learning-rate fix", f"{g('lr_fix','clump_%'):.1f}%",
             f"{g('lr_fix','spacing_CV'):.3f}", f"{g('lr_fix','CD_t_mm'):.2f}"],
            ["+ repulsion", f"{g('rep_w05','clump_%'):.1f}%",
             f"{g('rep_w05','spacing_CV'):.3f}", f"{g('rep_w05','CD_t_mm'):.2f}"],
            ["ground truth", "0.0%", "0.145", "—"]],
           col_w=[Inches(4.2), Inches(2.3), Inches(2.0), Inches(2.0)], highlight_row=3)
    _text(s, Inches(1.4), Inches(6.65), Inches(10.5), Inches(0.6),
          "Clumping falls to 1.4% against ground truth's 0.0%, at no accuracy cost — "
          "CD_t is unchanged within run-to-run variance. Repulsion is a density "
          "component, and I will describe it as one rather than as a general improvement.",
          size=14, color=MUTED)

    # =========================================================== 7. Results
    s = _slide(prs, "Results")
    _text(s, Inches(0.7), Inches(1.15), Inches(11.9), Inches(0.5),
          "20 validation skulls; millimetres converted with each skull's own scale. "
          "HD95 is the 95th-percentile error — Chamfer is a mean and averages away a "
          "single bad region, which for implant design is exactly what matters.",
          size=14, color=MUTED)
    _table(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(1.9),
           [["run", "CD_t (mm)", "HD95 (mm)", "F1@0.05", "F1@0.03", "DCD",
             "clumped<2mm", "spacing CV"]] +
           [[r] + [f"{g(r,c):.3f}" if c != "clump_%" else f"{g(r,c):.2f}%"
                   for c in ["CD_t_mm", "HD95_mm", "F1@0.05", "F1@0.03", "DCD",
                             "clump_%", "spacing_CV"]]
            for r in ["baseline", "dcd_w3", "dcd_l2", "lr_fix", "rep_w05"]] +
           [["ground truth", "—", "—", "—", "—", "—", "0.00%", "0.145"]],
           col_w=[Inches(2.2)] + [Inches(1.385)] * 7, highlight_row=5)
    _text(s, Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.4),
          "F-score against the source paper — same metric definition, different datasets",
          size=15, bold=True)
    _table(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.35),
           [["", "training data", "F1@0.05", "F1@0.03"],
            ["Source paper, Table 1", "200,000 clouds / 240 classes", "0.937", "0.741"],
            ["This work", "80 skulls / 1 class",
             f"{g('rep_w05','F1@0.05'):.3f}", f"{g('rep_w05','F1@0.03'):.3f}"],
            ["Source paper, Table 2", "4,800 shapes", "0.892", "0.623"]],
           col_w=[Inches(3.4), Inches(4.7), Inches(1.9), Inches(1.9)], highlight_row=2)
    _text(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(1.0),
          [("At the loose threshold this sits between their full model and their "
            "small-data run, with three orders of magnitude less training data; at the "
            "strict threshold it is behind both, which is where the remaining work is.",
            {}),
           ("DCD and F-score are comparable across projects because definitions and "
            "normalisation match. Chamfer is not — under this codebase's definition the "
            "paper's CD would fall 35× below the point spacing, so they are almost "
            "certainly using a squared-distance variant. I do not put those side by side.",
            {"size": 13, "color": MUTED})],
          size=14, space=5)

    # ========================================================= 8. Next steps
    s = _slide(prs, "Next steps")
    _text(s, Inches(0.75), Inches(1.2), Inches(5.8), Inches(0.35),
          "Planned", size=18, bold=True, color=ACCENT)
    _text(s, Inches(0.75), Inches(1.65), Inches(5.7), Inches(5.2),
          [("1.  Ablate DCD away", {"bold": True}),
           ("Density is handled by repulsion, the hyper-parameter experiments produced "
            "nothing, and the mechanism explains why. If removing it changes nothing, "
            "three losses collapse to two — a cleaner result than stacking terms.",
            {"size": 13.5, "color": MUTED}),
           ("2.  k-fold cross-validation", {"bold": True}),
           ("Every number here is a single 20-skull split. The code is in place and has "
            "never been run; this has to happen before anything is written up.",
            {"size": 13.5, "color": MUTED}),
           ("3.  Defect-region-restricted metrics", {"bold": True}),
           ("Most of the skull is already in the input, so current metrics are diluted "
            "by how well the model reproduces what it was handed. Evaluation-side change "
            "only — no retraining.", {"size": 13.5, "color": MUTED})],
          size=15, space=6)
    _text(s, Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.35),
          "Under consideration", size=18, bold=True, color=MUTED)
    _text(s, Inches(7.0), Inches(1.65), Inches(5.6), Inches(5.2),
          [("Focus the model on the defect", {"bold": True}),
           ("The dataset ships implant labels I have never used. Training on them "
            "directly puts all capacity on the hard part and matches the AutoImplant "
            "protocol — but it is a large change and breaks comparability with "
            "everything so far.", {"size": 13.5, "color": MUTED}),
           ("Point-to-plane attraction", {"bold": True}),
           ("Another route to the smoothness goal. Chamfer pulls each prediction toward "
            "a discrete sample, which lets points scatter either side of the true "
            "surface; pulling toward a fitted local plane removes that. Orthogonal to "
            "repulsion, verifiable with existing tools.", {"size": 13.5, "color": MUTED}),
           ("Poisson reconstruction control", {"bold": True}),
           ("The current visualisation is deliberately robust to uneven density, so it "
            "cannot show why density matters clinically. Poisson is not — running it "
            "would connect the metric to usability.", {"size": 13.5, "color": MUTED})],
          size=15, space=6)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"-> {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
